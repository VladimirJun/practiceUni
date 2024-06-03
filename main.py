import os
import tkinter as tk
from tkinter import messagebox
from tkinter import ttk

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


# Функция для добавления изобр ажения на слайд как фбэкграунд
def set_background(slide, image_path, prs):
    if not os.path.exists(image_path):
        messagebox.showerror("Ошибка", f"Файл {image_path} не найден.")
        return

    slide_width = prs.slide_width
    slide_height = prs.slide_height
    slide.shapes.add_picture(image_path, 0, 0, width=slide_width, height=slide_height)


# Функция для создания презентации
def create_presentation():
    game_class = class_level_var.get()
    if game_class == 1:
        class_column = 'Сложность для 1-4 класса'
    elif game_class == 5:
        class_column = 'Сложность для 5-7 класса'
    elif game_class == 8:
        class_column = 'Сложность для 8-11 класса'
    else:
        messagebox.showerror("Ошибка", "Выберите класс.")
        return

    try:
        num_easy_questions = int(num_questions_easy_var.get()) if num_questions_easy_var.get() else 0
        num_medium_questions = int(num_questions_mid_var.get()) if num_questions_mid_var.get() else 0
        num_hard_questions = int(num_questions_hard_var.get()) if num_questions_hard_var.get() else 0
    except ValueError:
        messagebox.showerror("Ошибка", "Введите корректное количество вопросов.")
        return

    excel_file_path = 'database.xlsx'
    if not os.path.exists(excel_file_path):
        messagebox.showerror("Ошибка", "Файл database.xlsx не найден.")
        return

    df = pd.read_excel(excel_file_path)

    filtered_easy = df[(df[class_column] < 5) & (df[class_column] > 0)]
    filtered_medium = df[(df[class_column] >= 5) & (df[class_column] < 8)]
    filtered_hard = df[df[class_column] >= 8]

    # проверка на достаточность вопросов
    if len(filtered_easy) < num_easy_questions:
        messagebox.showerror("Ошибка", "Не хватает легких вопросов.")
        return
    if len(filtered_medium) < num_medium_questions:
        messagebox.showerror("Ошибка", "Не хватает средних вопросов.")
        return
    if len(filtered_hard) < num_hard_questions:
        messagebox.showerror("Ошибка", "Не хватает сложных вопросов.")
        return

    # выбираем случайные вопросы
    selected_easy = filtered_easy.sample(n=num_easy_questions)
    selected_medium = filtered_medium.sample(n=num_medium_questions)
    selected_hard = filtered_hard.sample(n=num_hard_questions)

    # объединяем все выбранные вопросы
    selected_questions = pd.concat([selected_easy, selected_medium, selected_hard])
    # перемешиваем все вопросы и сбрасываем старые индексы
    selected_questions = selected_questions.sample(frac=1).reset_index(drop=True)

    # Создаем новую презентацию с измененными размерами слайда
    prs = Presentation()
    prs.slide_width = Inches(16)  # Ширина слайда
    prs.slide_height = Inches(9)  # Высота слайда

    # Добавляем первый слайд с фоном
    first_slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_background(first_slide, 'first.jpg', prs)

    for index, row in selected_questions.iterrows():
        game_title = row['Название игры']
        question = row['Вопрос']
        answer = row['Ответ']

        # Создаем слайд с вопросом
        question_slide = prs.slides.add_slide(prs.slide_layouts[1])
        question_title = question_slide.shapes.title
        question_content = question_slide.placeholders[1]
        question_title.text = f"ВОПРОС {index + 1}"
        if (index == len(selected_questions) - 1):
            question_title.text = f"ВОПРОС {index + 1} (ПОСЛЕДНИЙ)"
        question_content.text = f"{question}"

# Создаем слайд с ответом
        answer_slide = prs.slides.add_slide(prs.slide_layouts[5])
        set_background(answer_slide, 'answer.jpg', prs)
        textbox = answer_slide.shapes.add_textbox(Inches(3), Inches(3), Inches(14), Inches(7))
        text_frame = textbox.text_frame
        p = text_frame.add_paragraph()
        p.font.size = Pt(45)
        p.text = f"\n\t\t\t{answer}"
        p.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

    prs.save('output_presentation.pptx')
    messagebox.showinfo("Успех", "Презентация успешно создана!")

# Создание окна приложения
root = tk.Tk()
root.title("Project Practice")
root.geometry("500x400")
root.configure(bg="#F0F0F0")

# Надпись "Добро пожаловать"
welcome_label = tk.Label(root, text="Добро пожаловать", bg="#F0F0F0", font=("Arial", 20))
welcome_label.pack(pady=10)

# Выбор типа игры
game_type_frame = tk.Frame(root, bg="#F0F0F0")
game_type_frame.pack(pady=10)

game_type_label = tk.Label(game_type_frame, text="Выберите тип игры:", bg="#F0F0F0", font=("Arial", 12))
game_type_label.grid(row=0, column=0, padx=10)

game_type_var = tk.StringVar(root)
game_type_combobox = ttk.Combobox(game_type_frame, textvariable=game_type_var, values=["ЧГК"])
game_type_combobox.grid(row=0, column=1, padx=10)

# Поле ввода для количества лёгких вопросов
num_questions_label = tk.Label(root, text="Введите количество лёгких вопросов:", bg="#F0F0F0", font=("Arial", 12))
num_questions_label.pack()

num_questions_easy_var = tk.StringVar()
num_questions_entry = tk.Entry(root, textvariable=num_questions_easy_var)
num_questions_entry.pack()

# Поле ввода для количества средних вопросов
num_questions_label2 = tk.Label(root, text="Введите количество средних вопросов:", bg="#F0F0F0", font=("Arial", 12))
num_questions_label2.pack()

num_questions_mid_var = tk.StringVar()
num_questions_entry2 = tk.Entry(root, textvariable=num_questions_mid_var)
num_questions_entry2.pack()

# Поле ввода для количества сложных вопросов
num_questions_label3 = tk.Label(root, text="Введите количество сложных вопросов:", bg="#F0F0F0", font=("Arial", 12))
num_questions_label3.pack()

num_questions_hard_var = tk.StringVar()
num_questions_entry3 = tk.Entry(root, textvariable=num_questions_hard_var)
num_questions_entry3.pack()

# Выбор класса
class_level_frame = tk.Frame(root, bg="#F0F0F0")
class_level_frame.pack(pady=10)

class_level_label = tk.Label(class_level_frame, text="Выберите класс:", bg="#F0F0F0", font=("Arial", 12))
class_level_label.grid(row=0, column=0, padx=10)

class_level_var = tk.IntVar()

class_level_1_4 = tk.Radiobutton(class_level_frame, text="1-4", variable=class_level_var, value=1, bg="#F0F0F0",
                                 font=("Arial", 12))
class_level_1_4.grid(row=0, column=1)

class_level_5_7 = tk.Radiobutton(class_level_frame, text="5-7", variable=class_level_var, value=5, bg="#F0F0F0",
                                 font=("Arial", 12))
class_level_5_7.grid(row=0, column=2)

class_level_8_11 = tk.Radiobutton(class_level_frame, text="8-11", variable=class_level_var, value=8, bg="#F0F0F0",
                                  font=("Arial", 12))
class_level_8_11.grid(row=0, column=3)

# Кнопка "Создать презентацию"
create_button = tk.Button(root, text="Создать презентацию", bg="green", fg="white", font=("Arial", 12),
                          command=create_presentation)
create_button.pack(pady=20, fill="x")

root.mainloop()